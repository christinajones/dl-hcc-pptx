# script to automate generation of powerpoint file with slices from each volume
# creates a title slide and one slide per volume
# each volume slide displays thirty-five slices arranged in a grid

# uses python-pptx: https://python-pptx.readthedocs.io/en/latest/

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)

SLD_LAYOUT_BLANK=6
slide_layout=prs.slide_layouts[SLD_LAYOUT_BLANK]

# title slide # 
slide=prs.slides.add_slide(slide_layout)
background=slide.background
fill=background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)
left = Inches(0)
top = Inches(2.25)
width=Inches(13.33)
height=Inches(1.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
deck_title=tf.add_paragraph()
slide_title= 'Title Here' # add title for title slide here
deck_title.text = slide_title
deck_title.font.name='Garamond'
deck_title.font.size=Pt(55)
deck_title.alignment=PP_ALIGN.CENTER
deck_title.font.color.rgb= RGBColor(0xFF, 0xFF, 0xFF)

left = Inches(0)
top = Inches(4)
width=Inches(13.33)
height=Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
pf = txBox.text_frame
slide_subtitle= 'Subtitle Here' # add subtitle here
deck_subtitle=pf.add_paragraph()
deck_subtitle.text = slide_subtitle
deck_subtitle.font.name='Garamond'
deck_subtitle.font.size=Pt(30)
deck_subtitle.alignment=PP_ALIGN.CENTER
deck_subtitle.font.color.rgb=RGBColor(0xFF, 0xFF, 0xFF)

left = Inches(0)
top = Inches(5.25)
width=Inches(13.33)
height=Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
slide_date= 'Date Here' # add date here
deck_date=tf.add_paragraph()
deck_date.text = slide_date
deck_date.font.name='Garamond'
deck_date.font.size=Pt(18)
deck_date.alignment=PP_ALIGN.CENTER
deck_date.font.color.rgb=RGBColor(0xFF, 0xFF, 0xFF)

# volume slides #
for i in range(0, 131):
	slide=prs.slides.add_slide(slide_layout)
	background=slide.background
    	fill=background.fill
    	fill.solid()
    	fill.fore_color.rgb = RGBColor(0, 0, 0)
	IMG_HEIGHT=Inches(1.4) #image height value
   	IMG_WIDTH=Inches(1.74) #image width value
	j=0 
	 for j in range(35):
	 	# note: these image names might change depending upon how images have been named
		if j < 10: image_name='image_0'+str(j)+'.png' 
        	elif j >= 10: image_name='image_'+str(j)+'.png'
        	# note: this image path might change depending upon where images have been saved
       		img_path="screenshots/"+str(i)+"/"+image_name 
		
		#determining top (vertical placement) value for picture
		if j>=0 and j < 7: top =Inches(0.4)
     		elif j>=7 and j < 14: top = Inches(1.8)
        	elif j>=14 and j<21: top = Inches(3.2)
        	elif j>=21 and j < 28: top = Inches(4.56)
        	elif j>=28 and j < 35: top = Inches(6)
		
		#determining left (horizontal placement) value for picture
	 	if (j%7)==0: left=Inches(0.5)
        	elif (j%7)==1: left=Inches(2.24)
        	elif (j%7)==2: left=Inches(3.98)
        	elif (j%7)==3: left=Inches(5.72)
        	elif (j%7)==4: left=Inches(7.46)
        	elif (j%7)==5: left=Inches(9.2)
        	elif (j%7)==6: left=Inches(10.94)
	
	picture=slide.shapes.add_picture(img_path, left, top, width=IMG_WIDTH, height=IMG_HEIGHT)
	left = Inches(0)
    top = Inches(-0.25)
    width=Inches(1.5)
    height=Inches(0.5)
	txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    slide_volume='Volume ' + str(i)
    slide_title=tf.add_paragraph()
    slide_title.text=slide_volume
    slide_title.font.name='Garamond'
    slide_title.font.size=Pt(18)
    slide_title.font.color.rgb=RGBColor(0xFF, 0xFF, 0xFF)

prs.save('filename.pptx') # enter desired file name for generated powerpoint document